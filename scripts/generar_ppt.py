"""
Script para generar la Presentación Ejecutiva en formato PowerPoint (.pptx)
para el Proyecto Biblioteca Normativa Circulares DDU.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colores del Tema Ejecutivo
NAVY_DARK = RGBColor(15, 23, 42)      # #0F172A
NAVY_SURFACE = RGBColor(30, 41, 59)   # #1E293B
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38BDF8
TEXT_PRIMARY = RGBColor(248, 250, 252)# #F8FAFC
TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8
WHITE = RGBColor(255, 255, 255)
BG_LIGHT = RGBColor(241, 245, 249)    # #F1F5F9
CARD_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(203, 213, 225)
TEXT_DARK = RGBColor(15, 23, 42)
RED_ACCENT = RGBColor(239, 68, 68)
GREEN_ACCENT = RGBColor(34, 197, 94)

def aplicar_fondo_oscuro(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_DARK

def agregar_encabezado_diapositiva(slide, titulo_texto, subtitulo_texto):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = titulo_texto
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY_DARK
    p.font.name = "Segoe UI"
    
    p2 = tf.add_paragraph()
    p2.text = subtitulo_texto
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(100, 116, 139)
    p2.font.name = "Segoe UI"

def construir_presentacion(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # DIAPOSITIVA 1: Portada (Estilo Dark Navy)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    aplicar_fondo_oscuro(slide1)

    tx_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "PROYECTO BIBLIOTECA NORMATIVA CIRCULARES DDU"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Segoe UI"
    p.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "Estructuración Semántica (Akoma Ntoso XML & Grafos RDF)"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_PRIMARY
    p2.font.name = "Segoe UI"
    p2.space_after = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Potenciando la precisión y calidad de respuestas del Agente AI Legal para las DOM"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Segoe UI"

    # -------------------------------------------------------------
    # DIAPOSITIVA 2: El Desafío (RAG con PDFs vs Akoma+RDF)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    agregar_encabezado_diapositiva(
        slide2,
        "El Desafío Institucional y el Agente 'Biblioteca Normativa'",
        "De la Ingesta Tradicional RAG no estructurada a la Estructuración Semántica Atómica"
    )

    # Texto Explicativo de Contexto
    tx_ctx = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.8))
    tf_ctx = tx_ctx.text_frame
    tf_ctx.word_wrap = True
    p_ctx = tf_ctx.paragraphs[0]
    p_ctx.text = "El aplicativo 'Biblioteca Normativa' opera mediante arquitectura RAG (Retrieval-Augmented Generation). No obstante, la ingesta directa de PDFs planos no estructurados producía respuestas de mala calidad y falta de precisión jurídica. La incorporación del estándar Akoma Ntoso XML y su Grafo RDF provee una estructura clara y relaciones definidas para elevar radicalmente la precisión."
    p_ctx.font.size = Pt(13)
    p_ctx.font.color.rgb = TEXT_DARK
    p_ctx.font.name = "Segoe UI"

    # Tabla Comparativa
    rows, cols = 4, 3
    left, top, width, height = Inches(0.8), Inches(2.9), Inches(11.733), Inches(3.8)
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.6)
    table.columns[2].width = Inches(4.633)

    headers = ["Dimensión", "Antes (Documentos No Estructurados con RAG)", "Con la Solución (Akoma Ntoso + Grafo RDF)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"

    data = [
        ("Ingesta de Datos", "PDFs planos procesados en bloques de texto arbitrarios sin jerarquía.", "Estructura atómica jerárquica (secciones/párrafos) y grafo de relaciones."),
        ("Calidad de Respuesta", "Mala calidad general en las respuestas y riesgo de fragmentación de contexto.", "Respuestas exactas, fundamentadas y con cita legal atómica por párrafo."),
        ("Fundamentación Legal", "Dificultad para precisar numerales, incisos y notas explicativas al pie.", "Identificación precisa de numerales, secciones, notas al pie y descriptores.")
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            if col_idx == 0:
                cell.fill.fore_color.rgb = BG_LIGHT
            elif col_idx == 1:
                cell.fill.fore_color.rgb = RGBColor(254, 242, 242)
            else:
                cell.fill.fore_color.rgb = RGBColor(240, 253, 244)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = "Segoe UI"
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY_DARK
            elif col_idx == 1:
                p.font.color.rgb = RGBColor(185, 28, 28)
            else:
                p.font.color.rgb = RGBColor(21, 128, 61)

    # -------------------------------------------------------------
    # DIAPOSITIVA 3: Arquitectura y Flujo End-to-End
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    agregar_encabezado_diapositiva(
        slide3,
        "Arquitectura de Procesamiento y Flujo General End-to-End",
        "Tubería de transformación dinámica de 5 etapas desde el PDF fuente hasta la capa de datos semánticos"
    )

    pasos = [
        ("1. PDF DDU", "Documento fuente original publicado por la División de Desarrollo Urbano."),
        ("2. Extractores ETL", "11 extractores modulares independientes que procesan el texto dinámicamente."),
        ("3. CSV Dominio", "Capa intermedia de datos tabulares planos de dominio de la circular."),
        ("4. Akoma XML BCN", "Estructura jurídica atómica conforme al estándar oficial BCN v2.0."),
        ("5. Grafo RDF", "Publicación en formato Turtle para navegación en la web de datos abiertos.")
    ]

    box_w = Inches(2.15)
    box_h = Inches(4.2)
    gap = Inches(0.24)
    start_x = Inches(0.8)
    start_y = Inches(2.2)

    for i, (titulo_paso, desc_paso) in enumerate(pasos):
        x = start_x + i * (box_w + gap)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ACCENT_BLUE if i in (3, 4) else BORDER_COLOR
        shape.line.width = Pt(2 if i in (3, 4) else 1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = titulo_paso
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(12)

        p2 = tf.add_paragraph()
        p2.text = desc_paso
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(71, 85, 105)
        p2.font.name = "Segoe UI"

    # -------------------------------------------------------------
    # DIAPOSITIVA 4: El Valor de la Estructuración Semántica
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    agregar_encabezado_diapositiva(
        slide4,
        "El Valor de la Estructuración Semántica (Akoma Ntoso + RDF)",
        "Dos pilares fundamentales para potenciar las capacidades de respuesta del Agente AI"
    )

    card_w = Inches(5.6)
    card_h = Inches(4.5)

    # Pilar 1
    c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), card_w, card_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = WHITE
    c1.line.color.rgb = NAVY_DARK
    c1.line.width = Pt(1.5)

    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "📐 Pilar 1: Estructura Jurídica Estándar"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY_DARK
    p.font.name = "Segoe UI"
    p.space_after = Pt(6)

    p_sub = tf1.add_paragraph()
    p_sub.text = "(Akoma Ntoso BCN)"
    p_sub.font.bold = True
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = ACCENT_BLUE
    p_sub.font.name = "Segoe UI"
    p_sub.space_after = Pt(16)

    p_body1 = tf1.add_paragraph()
    p_body1.text = "• Permite al Agente citar numerales, párrafos, notas al pie y descriptores exactos.\n\n• Otorga una jerarquía clara de secciones y párrafos atómicos etiquetados con identificadores únicos."
    p_body1.font.size = Pt(13)
    p_body1.font.color.rgb = RGBColor(51, 65, 85)
    p_body1.font.name = "Segoe UI"

    # Pilar 2
    c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), card_w, card_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = WHITE
    c2.line.color.rgb = NAVY_DARK
    c2.line.width = Pt(1.5)

    tf2 = c2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = Inches(0.3)

    p = tf2.paragraphs[0]
    p.text = "🕸️ Pilar 2: Red de Conocimiento Legal"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY_DARK
    p.font.name = "Segoe UI"
    p.space_after = Pt(6)

    p_sub2 = tf2.add_paragraph()
    p_sub2.text = "(Grafo RDF Turtle)"
    p_sub2.font.bold = True
    p_sub2.font.size = Pt(14)
    p_sub2.font.color.rgb = ACCENT_BLUE
    p_sub2.font.name = "Segoe UI"
    p_sub2.space_after = Pt(16)

    p_body2 = tf2.add_paragraph()
    p_body2.text = "• Permite al Agente navegar las relaciones entre circulares.\n\n• Conecta semánticamente qué circulares complementan o modifican a otras dentro del ecosistema normativo."
    p_body2.font.size = Pt(13)
    p_body2.font.color.rgb = RGBColor(51, 65, 85)
    p_body2.font.name = "Segoe UI"

    # -------------------------------------------------------------
    # DIAPOSITIVA 5: Próximos Pasos e Hoja de Ruta
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    agregar_encabezado_diapositiva(
        slide5,
        "Próximos Pasos e Hoja de Ruta de Integración",
        "Plan de acción para el acoplamiento técnico al aplicativo y pruebas de rendimiento del Agente"
    )

    fase_w = Inches(11.733)
    fase_h = Inches(2.1)

    # Fase 1
    sh_f1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), fase_w, fase_h)
    sh_f1.fill.solid()
    sh_f1.fill.fore_color.rgb = RGBColor(240, 249, 255)
    sh_f1.line.color.rgb = ACCENT_BLUE
    sh_f1.line.width = Pt(2)

    tf_f1 = sh_f1.text_frame
    tf_f1.word_wrap = True
    tf_f1.margin_left = tf_f1.margin_right = tf_f1.margin_top = Inches(0.25)

    pf1 = tf_f1.paragraphs[0]
    pf1.text = "Fase 1: Integración Inicial y Pruebas de Rendimiento (Punto de Partida Baseline)"
    pf1.font.bold = True
    pf1.font.size = Pt(16)
    pf1.font.color.rgb = NAVY_DARK
    pf1.font.name = "Segoe UI"
    pf1.space_after = Pt(8)

    pf1_b = tf_f1.add_paragraph()
    pf1_b.text = "El equipo técnico acopla la nueva estructura (Akoma Ntoso XML y Grafo RDF actualmente construidos) al aplicativo 'Biblioteca Normativa' como punto de partida ('un desde'). Posteriormente, el equipo interno de negocio genera pruebas de rendimiento para ver cómo responde el Agente ante esta nueva ingesta de información asociada a las circulares."
    pf1_b.font.size = Pt(12)
    pf1_b.font.color.rgb = RGBColor(30, 41, 59)
    pf1_b.font.name = "Segoe UI"

    # Fase 2
    sh_f2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), fase_w, fase_h)
    sh_f2.fill.solid()
    sh_f2.fill.fore_color.rgb = WHITE
    sh_f2.line.color.rgb = BORDER_COLOR
    sh_f2.line.width = Pt(1.5)

    tf_f2 = sh_f2.text_frame
    tf_f2.word_wrap = True
    tf_f2.margin_left = tf_f2.margin_right = tf_f2.margin_top = Inches(0.25)

    pf2 = tf_f2.paragraphs[0]
    pf2.text = "Fase 2: Escalamiento y Complejidad Estructural"
    pf2.font.bold = True
    pf2.font.size = Pt(16)
    pf2.font.color.rgb = NAVY_DARK
    pf2.font.name = "Segoe UI"
    pf2.space_after = Pt(8)

    pf2_b = tf_f2.add_paragraph()
    pf2_b.text = "Profundizar y escalar la extracción procesando nuevas circulares DDU que posean variaciones estructurales más complejas, consolidando la cobertura completa del catálogo normativo y ajustando el orquestador ante nuevos patrones."
    pf2_b.font.size = Pt(12)
    pf2_b.font.color.rgb = RGBColor(71, 85, 105)
    pf2_b.font.name = "Segoe UI"

    # Guardar Presentación
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentación PPT generada exitosamente en: {output_path}")

if __name__ == "__main__":
    output_ppt = Path("salidas_ppt/Presentacion_Ejecutiva_Biblioteca_Normativa_DDU.pptx")
    construir_presentacion(output_ppt)

import sys
import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx not found. Installing...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Análisis de Estándares y Normas"
    subtitle.text = "API Ley Chile de la BCN\nPropuesta de Plan de Acción"

    # Slide 1
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "1. Etapa del Proceso de Datos"
    tf = body_shape.text_frame
    tf.text = "Orientación de los Estándares de la API:"
    p = tf.add_paragraph()
    p.text = "Etapa de distribución, acceso, interoperabilidad y consumo de datos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "No aborda la creación (redacción o firma de leyes), sino la exposición desde el repositorio digital."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Gestión Jerárquica: Estructura en árbol de la norma (servicio/9)."
    p.level = 1

    # Slide 2
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "2. Sistemas Tecnológicos"
    tf = body_shape.text_frame
    tf.text = "¿Apuntan a generar data para sistemas tecnológicos?"
    p = tf.add_paragraph()
    p.text = "Sí, su objetivo es absoluto y principal."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Protocolo de Sistemas: Es una API RESTful documentada bajo OpenAPI 3.1.0, ideal para comunicación Máquina a Máquina (M2M)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Seguridad: Uso de parámetro de autenticación ('secret') en los endpoints, asegurando integraciones formales."
    p.level = 1

    # Slide 3
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "3. Datos para Inteligencia Artificial"
    tf = body_shape.text_frame
    tf.text = "¿Ayudan a levantar datos para IA?"
    p = tf.add_paragraph()
    p.text = "Sí, son extremadamente valiosos y útiles (evitan el OCR 'en bruto')."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Granularidad: Fragmentos exactos de ley por 'idParte' (artículos)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Grafos de Conocimiento: Relaciones explícitas de MODIFICACIÓN y CONCORDANCIA (servicio/846)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Vigencia temporal explícita para evitar alucinaciones normativas."
    p.level = 1

    # Slide 4
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "4. Formato de los Datos"
    tf = body_shape.text_frame
    tf.text = "La API utiliza formatos de estándar mundial:"
    p = tf.add_paragraph()
    p.text = "XML (application/xml): Formato predominante, excelente para representar jerarquías (Títulos, Párrafos, Artículos)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "JSON (application/json): Estándar moderno web, ideal para integración directa con IA y lenguajes modernos. (Ej. servicio/7.2)."
    p.level = 1

    # Slide 5
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "5. Conclusión Estratégica"
    tf = body_shape.text_frame
    tf.text = "¿Crear plan de acción con Osvaldo?"
    p = tf.add_paragraph()
    p.text = "Sí. Cambia radicalmente la eficiencia y disminuye el margen de error."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Delega el mantenimiento normativo a la BCN."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Permite enfocar los esfuerzos en valor agregado: sistemas de IA, RAG y búsquedas semánticas."
    p.level = 1

    # Slide 6
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Plan de Acción (Propuesta)"
    tf = body_shape.text_frame
    tf.text = "Fases de Implementación:"
    p = tf.add_paragraph()
    p.text = "Etapa 1: Gestión de Accesos (Obtener 'secret' de la BCN)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Etapa 2: Pruebas de Concepto y Mapeo (JSON/XML a modelo propio)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Etapa 3: Integración Crítica (Sincronización masiva)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Etapa 4: Automatización (Consulta de cambios diarios)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Etapa 5: Ingesta en Sistemas de IA (Alimentar bases vectoriales)."
    p.level = 1

    ppt_path = r"C:\Users\preusc\.gemini\antigravity-cli\brain\26d7dfbd-4460-4b72-8345-68b4afcf6448\Presentacion_API_Ley_Chile.pptx"
    prs.save(ppt_path)
    print(f"PPT generado con éxito en: {ppt_path}")

if __name__ == '__main__':
    create_presentation()
